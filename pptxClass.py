from pptx import Presentation
from pptx import util
import os
import datetime
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptxClass import *

class PPTX:
    def __init__(self, fileName,folder,template=None):
        self.fn = fileName
        self.folder = folder
        self.template = template

        if not os.path.exists(folder):
            os.makedirs(folder)

        self.prs = Presentation(self.template)

        self.wd = self.prs.slide_width
        self.hg = self.prs.slide_height

    def addWeek(self,week):
        blank_slide_layout=self.prs.slide_layouts[0]
        slide=self.prs.slides.add_slide(blank_slide_layout)
        for i,ph in enumerate(slide.placeholders):
            ph.text=week[i]

        #version = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        print('ciaone')
        self.savePres()
        return
    '''
    def addTable(self,title,footer,matrix,header=None,index=None,rnd=2):
        blank_slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        slide.placeholders[10].text = title
        slide.placeholders[11].text = footer

        rows,cols = matrix.shape
        if header != None:
            rows +=1
        if index != None:
            cols +=1
        x = int(self.wd/20)
        y = int(self.hg/10)
        cx = util.Cm(20)
        cy = util.Cm(4)

        shape = slide.shapes.add_table(rows,cols, x, y, cx, cy)
        table = shape.table

        si = 0
        sj = 0

        if index is not None:
            sj = 1
        if header is not None:
            si = 1


        #these three lines are needed to avoid having table(0,0) white in case of index && header
        cell = table.cell(0,0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220,220,220)

        if index is not None:
            for i in range(len(index)):
                cell=table.cell(i+si,0)
                cell.text = index[i]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220,220,220)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.bold = True

        if header is not None:
            for i in range(len(header)):
                cell=table.cell(0,i+sj)
                cell.text = header[i]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220,220,220)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.bold = True

        for i in range(si,rows):
            for j in range(sj,cols):
                cell=table.cell(i,j)
                cell.text = str(round(matrix[i-si][j-sj],rnd))
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255,255,255)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)

        self.savePres()
        return

    '''
    def savePres(self):
        self.prs.save(self.folder+'/'+self.fn+'.pptx')
        return




